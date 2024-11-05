import argparse
import os
import shutil
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from langchain_community.embeddings.ollama import OllamaEmbeddings
from langchain_community.vectorstores.chroma import Chroma
from langchain.prompts import ChatPromptTemplate
from langchain.schema.document import Document as LangchainDocument
from langchain_groq import ChatGroq
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import pytesseract
from PIL import Image
from io import BytesIO
from reportlab.pdfgen import canvas
import io 
import tabula
from textwrap import wrap
import warnings
from langchain_core.output_parsers import StrOutputParser
warnings.filterwarnings("ignore", category=FutureWarning, module='tabula.io')

#sostituire con le vostre directory
pytesseract.pytesseract.tesseract_cmd = 'C:/Program Files (x86)/Tesseract/tesseract.exe' 
DATA_PATH = "C:\\Users\\Rinaldi\\Desktop\\tirocinio\\file"
PARSED_DATA_PATH = "C:\\Users\\Rinaldi\\Desktop\\tirocinio\\temp_files"
CHROMA_PATH = "C:\\Users\\Rinaldi\\Desktop\\tirocinio\\chromadb"
PROMPT_TEMPLATE = """
Answer the question based only on the following context:

{context}

---

# Answer the question based on the above context: {question}
# """

def main():
    # Crea il parser degli argomenti
    parser = argparse.ArgumentParser()

    # Aggiungi l'argomento per resettare il database
    parser.add_argument("--reset", action="store_true", help="Reset the database.")

    # Parsing degli argomenti dalla riga di comando
    args = parser.parse_args()

    # Controlla se è richiesto il reset del database
    if args.reset:
        print("✨ Clearing Database")
        clear_database()

    # Carica i documenti, suddividi e aggiungi al database vettoriale
    documents = load_documents()
    chunks = split_documents(documents)
    add_to_chroma(chunks)

    # Inizia la chat
    chat()

#funzione chat
def chat():
    # Inizializza una lista per memorizzare le conversazioni precedenti
    conversation_history = []

    # Inizia il ciclo di chat
    while True:
        # Chiedi all'utente di inserire una domanda
        user_question = input("You😁: ")

        # Controlla se l'utente vuole uscire dalla chat
        if user_question.lower() in ["exit", "quit"]:
            print("Chat ended. Goodbye!")
            break

        # Crea il contesto per la query combinando le conversazioni precedenti
        context = "\n\n---\n\n".join(["You: " + q + "\nModel: " + r.content for q, r in conversation_history])

        # Chiedi una risposta al modello
        model_response = query_rag(user_question, context)

        # Stampa la risposta del modello
        # print(f"Model: {model_response}")

        # Aggiungi la domanda e la risposta alla memoria delle conversazioni
        conversation_history.append((user_question, model_response))

#funzione query
def query_rag(query_text: str, context: str):
    # Prepara il DB.
    embedding_function = get_embedding_function()
    db = Chroma(persist_directory=CHROMA_PATH, embedding_function=embedding_function)

    # Esplora DB.
    results = db.similarity_search_with_score(query_text, k=5)

    # Aggiungi il contesto al prompt
    context_text = context + "\n\n" + "\n\n---\n\n".join([doc.page_content for doc, _score in results])
    prompt_template = ChatPromptTemplate.from_template(PROMPT_TEMPLATE)
    prompt = prompt_template.format(context=context_text, question=query_text)

    # Esegui la query con il modello
    #model = Ollama(model="llama3:latest")
    #sostituite con la vostra key
    model = ChatGroq(model="llama3-70b-8192", api_key="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") 
    response_text = model.invoke(prompt)

    # Fonti
    sources = [doc.metadata.get("id", None) for doc, _score in results]
    formatted_response = f"Response🦙: {response_text}\nSources📖: {sources}\n\n"
    print(formatted_response)

    # Restituisci la risposta del modello
    return response_text

# Estensioni di file supportate
SUPPORTED_EXTENSIONS = ['.pdf', '.docx', '.pptx', '.txt']
#funzione per caricare i documenti

def load_documents():
    # Percorso della directory da cui caricare i documenti
    directory_path = DATA_PATH
    
    # Percorso per salvare i documenti di testo estratti
    parsed_data_path = PARSED_DATA_PATH
    
    # Verifica se la cartella esiste, altrimenti creala
    if not os.path.exists(parsed_data_path):
        os.makedirs(parsed_data_path)
    
    # Ottieni l'elenco di tutti i file nella directory
    files_in_directory = os.listdir(directory_path)
    
    # Filtra i file con estensioni supportate
    supported_files = [file for file in files_in_directory if os.path.splitext(file)[1] in SUPPORTED_EXTENSIONS and not file.startswith('~$')]
    
    # Lista per memorizzare tutti i documenti caricati
    documents = []
    
    # Carica ogni file
    for file in supported_files:
        # Costruisci il percorso completo del file
        file_path = os.path.join(directory_path, file)
        
        # Crea il nome del file PDF estratto
        file_name, file_extension = os.path.splitext(file)
        parsed_file_name = f"{file_name}_parsed.pdf"
        parsed_file_path = os.path.join(parsed_data_path, parsed_file_name)
        
        # Carica il documento a seconda del tipo di file
        text = ""  # Variabile per contenere il testo estratto
        
        if file_extension == '.pdf':
            # Carica file PDF e estrai testo
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                text += "\n".join([page.extract_text().strip() for page in reader.pages])
            
            # Utilizza tabula-py per estrarre le tabelle
            tables = tabula.read_pdf(file_path, pages='all')
            
           # Converti ogni tabella in linguaggio naturale
            for table in tables:
                prompt = "transform the following table into natural language including in every row transformed all the informations about what that value rappresents:" + table.to_string(index=False)
                model = ChatGroq(model="llama3-70b-8192", api_key="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")  
                chain = model | StrOutputParser()
                response_text = chain.invoke(prompt)
                # Rimuovi eventuali frasi aggiuntive non presenti nella tabella originale
                original_table_text = table.to_string(index=False)
                response_text = response_text.replace(original_table_text, "")
                text += f"{response_text}\n"
            
            
        elif file_extension == '.docx':
            # Carica file DOCX e estrai testo
            # Apri il file .docx
            doc = DocxDocument(file_path)
            
            # Estrai il testo dalle sezioni del documento
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Estrai tabelle e convertile in linguaggio naturale
            for table in doc.tables:
                df = pd.DataFrame([[cell.text for cell in row.cells] for row in table.rows])
                prompt = "transform the following table into natural language including in every row transformed all the informations about what that value rappresents:" + df.to_string(index=False)
                model = ChatGroq(model="llama3-70b-8192", api_key="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") 
                chain = model |StrOutputParser()
                response_text = chain.invoke(prompt)
                text += f"{response_text}"
            
            # Estrai le immagini e applica OCR
            for image_part in doc.part.rels.values():
                if image_part.reltype == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image":
                    image_data = image_part.target_part.blob
                    # Apri l'immagine come oggetto PIL
                    image = Image.open(io.BytesIO(image_data))
                    
                    # Esegui OCR sull'immagine usando pytesseract
                    ocr_text = pytesseract.image_to_string(image)
                    
                    # Aggiungi il testo estratto al documento
                    text += f"\n{ocr_text}\n"
            
        elif file_extension == '.pptx':
            # Carica file PPTX e estrai testo
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"
                    
                    # Estrarre tabelle e convertirle in linguaggio naturale
                    if shape.has_table:
                        table = shape.table
                        df = pd.DataFrame(
                            [[cell.text_frame.text for cell in row.cells] for row in table.rows]
                        )
                        prompt = "transform the following table into natural language including in every row transformed all the informations about what that value rappresents" + df.to_string(index=False)
                        model = ChatGroq(model="llama3-70b-8192", api_key="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") 
                        chain = model |StrOutputParser()
                        response_text = chain.invoke(prompt)
                        text += f"{response_text}"
                    
                    # Analizza le immagini nella presentazione
                    if shape.shape_type == 13:  # Tipo di forma per immagine
                        try:
                            # Esegui OCR sull'immagine
                            image = Image.open(BytesIO(shape.image.blob))
                            ocr_text = pytesseract.image_to_string(image)
                            text += f"\n{ocr_text}"
                        except Exception as e:
                            print(f"Error processing image in PPTX {file_path}: {e}")
            
        elif file_extension == '.txt':
            # Carica file TXT e estrai testo
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read().strip()
        
        save_to_pdf(text, parsed_file_path)
        # Aggiungi il documento al database
        documents.append(LangchainDocument(page_content=text, metadata={'source': file_path}))
    
    print(documents)
    # Restituisci la lista di documenti caricati
    return documents

def save_to_pdf(text, parsed_file_path):
    c = canvas.Canvas(parsed_file_path)
    lines = text.split("\n")
    y_position = 800  # Posizione iniziale verticale (coordinate PDF)

    # Aggiungi ogni linea al file PDF
    for line in lines:
        # Avvolgi il testo in righe di larghezza massima 70 caratteri (puoi regolare questo valore)
        wrapped_lines = wrap(line, width=70)

        # Scrivi ogni riga avvolta sul PDF
        for wrapped_line in wrapped_lines:
            c.drawString(50, y_position, wrapped_line)
            y_position -= 12  # Sposta la posizione verticale per la prossima riga

            # Se la posizione Y è troppo bassa, crea una nuova pagina
            if y_position < 100:
                c.showPage()
                y_position = 800

    # Termina il documento PDF
    c.save()
    print(f"✨ File parsed and saved: {parsed_file_path}")

#funzione per splittare i documenti
def split_documents(documents: list[Document]):

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=80,
        length_function=len,
        is_separator_regex=False,
    )
    return text_splitter.split_documents(documents)

#funzione di embedding
def get_embedding_function():
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    return embeddings

def add_to_chroma(chunks: list[LangchainDocument]):
    # Inizializza il database vettoriale Chroma
    db = Chroma(persist_directory=CHROMA_PATH, embedding_function=get_embedding_function())

    # Calcola gli ID delle pagine 
    chunks_with_ids = calculate_chunk_ids(chunks)

    # Ottieni gli ID dei documenti esistenti nel database
    existing_items = db.get(include=[])  # Gli ID sono sempre inclusi di default
    existing_ids = set(existing_items["ids"])
    print(f"Number of existing documents in DB: {len(existing_ids)}")

    # Filtra i documenti nuovi da aggiungere
    new_chunks = [chunk for chunk in chunks_with_ids if chunk.metadata["id"] not in existing_ids]

    if len(new_chunks) > 0:
        print(f"👉 Adding new documents: {len(new_chunks)}")
        new_chunk_ids = [chunk.metadata["id"] for chunk in new_chunks]
        db.add_documents(new_chunks, ids=new_chunk_ids)
        # Non è più necessario chiamare esplicitamente db.persist(), poiché Chroma gestisce la persistenza automaticamente.
        print("✅ New documents added.")
    else:
        print("✅ No new documents to add.")



def calculate_chunk_ids(chunks):
    # Pagina sorgente : Numero di pagina : Indice chunk
    last_page_id = None
    current_chunk_index = 0

    for chunk in chunks:
        source = chunk.metadata.get("source")
        page = chunk.metadata.get("page")
        current_page_id = f"{source}:{page}"

        if current_page_id == last_page_id:
            current_chunk_index += 1
        else:
            current_chunk_index = 0

        # Calcula il chunk ID.
        chunk_id = f"{current_page_id}:{current_chunk_index}"
        last_page_id = current_page_id

        # Aggiungilo ai metadati della pagina.
        chunk.metadata["id"] = chunk_id

    return chunks

def clear_database():
    if os.path.exists(CHROMA_PATH):
        shutil.rmtree(CHROMA_PATH)


if __name__ == "__main__":
    main()
